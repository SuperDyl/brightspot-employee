"""
Random utility functions used by the professor package
"""
from pptx import Presentation

import professor

from bs4 import BeautifulSoup
from typing import Iterable
import requests

BLANK_LAYOUT = Presentation().slide_layouts[6]
NAME_SUFFIXES = ['jr.', 'iii', 'sr.']


def split_name(full_name: str, name_suffixes: Iterable[str] = None) -> (str, str):
    """Return the estimated first and last name as a tuple"""
    if name_suffixes is None:
        name_suffixes = NAME_SUFFIXES
    full_split_name = full_name.split(' ')
    if full_split_name[-1].lower() in name_suffixes:
        *first, last = full_split_name[:-1]
        return ' '.join(first), ' '.join((last, full_split_name[-1]))
    # else
    *first, last = full_split_name
    return ' '.join(first), last


def tag_iterator(url, *args, **kwargs):
    if not args:
        args = 'div',
    if not kwargs:
        kwargs = {'class_': 'ListVerticalImage-items-item'}

    with requests.get(url) as request:
        html_data = request.text
    bs = BeautifulSoup(html_data, 'html.parser')
    return bs.find_all(*args, **kwargs)


def remove_prefix(string: str, pref: str):
    if pref == string[0:len(pref)]:
        return string[len(pref):-1]
    return string


def call_each(funcs: iter, *args, **kwargs):
    return (x(*args, **kwargs) for x in funcs)


# def print_bad_formatted_rooms():
#     update_professor_file()
#     old_p = pull_professor_data_old()
#     professors = parse_professor_data()
#     for old, new in zip(old_p, professors):
#         if old.room != (str(new.room)[3:]):
#             print(old)
#     print("Done!")
#     # print(*(f'{n.room}, {o.room}' for o, n in zip(old, professors)), sep='\n')


# def pull_professor_data_old(url="https://religion.byu.edu/directory"):
#     with urlopen(url) as request:
#         html_data = request.read().decode("utf-8")
#     bs = BeautifulSoup(html_data, 'html.parser')
#
#     professors = []
#     for tag in bs.find_all('div', class_='PromoVerticalImage-content'):
#         room = ''
#         try:
#             room = tag.find_all('p')[0].contents[0].string.strip()
#             # room = RoomAddress.from_string(room)
#         except IndexError:
#             pass
#         name = tag.find_all(class_='PromoVerticalImage-title promo-title')[0].find('a').string
#         professors.append(professor.Professor(name, room))
#
#     return professors


# def update_professor_file(file_name='professors.txt', professors=None):
#     if not professors:
#         professors = pull_professor_data(url="https://religion.byu.edu/directory")
#     new_prof_dict = dict(zip((x.name for x in professors), professors))
#     old_professors = parse_professor_data(file_name)
#     old_prof_dict = dict(zip((x.name for x in old_professors), old_professors))
#     old_prof_dict.update(new_prof_dict)
#     for name, prof in old_prof_dict.items():
#         prof.name = name
#
#     with open(file_name, 'w') as file:
#         file.write(('\n'.join((str(p) for p in old_prof_dict.values()))))
#
#     return old_prof_dict.values()
#
#
# def parse_professor_data(file_name='professors.txt'):
#     with open(file_name, 'r') as file:
#         return [professor.Professor.from_string(line) for line in file]


# professors = update_professor_file()
# os.listdir()
# print_bad_formatted_rooms()

# for n in names:
#     add_slide(n)
# prs.save("Testing2.pptx")

# professors = pull_professor_data()
# print(*professors, sep='\n')
#
# rooms = RoomAddress.from_string_iter(rooms)
# print(len(names), len(rooms))
# print(*(f'{name:25}={str(room)}' for name, room in zip(names, rooms)), sep='\n')
#
# names, rooms = pull_professor_data()
# print(*(r for r in rooms), sep='\n')
# print(*(r for r in RoomAddress.from_string_iter(rooms)), sep='\n')
# rooms = RoomAddress.from_string_iter(rooms)
# print(len(names), len(rooms))
# print(*(f'{name:25}={str(room)}' for name, room in zip(names, rooms)), sep='\n')
#
# names, rooms = pull_professor_data()
# print(*(f'{n:35} = {r}' for n, r in zip(names, rooms)), sep='\n')
# print(*(r for r in rooms), sep='\n')
#
# add_slide("John Doe")
# prs.save("Testing.pptx")
